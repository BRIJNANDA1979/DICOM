# -*- coding: utf-8 -*-
"""
Created on Tue Jul 23 16:16:08 2024

@author: BRIJB
"""

# -*- coding: utf-8 -*-
"""
Created on Mon Jul 22 14:18:14 2024

@author: BRIJB
"""
import os
from pptx import Presentation
from pptx.util import Inches,Pt
from datetime import datetime
from pptx.enum.text import PP_ALIGN
import json
import csv
import numpy as np
import sys
import re

import time
def list_level_subdirectories(start_dir):
    level = 0
    third_level_dirs = []

    # Use os.walk to traverse the directory
    for root, dirs, files in os.walk(start_dir):
        # Calculate the current level based on the depth of root
        current_level = root.replace(start_dir, '').count(os.sep)
        
        # Check if it's the third level
        if current_level == 3:
            for d in dirs:
                third_level_dirs.append(os.path.join(root, d))
        elif current_level > 3:
            break  # No need to go deeper once we reach beyond the third level

    return third_level_dirs

def list_level_third_subdirectories(start_dir):
    level = 0
    third_level_dirs = []

    # Use os.walk to traverse the directory
    for root, dirs, files in os.walk(start_dir):
        # Calculate the current level based on the depth of root
        current_level = root.replace(start_dir, '').count(os.sep)
        
        # Check if it's the third level
        if current_level == 2:
            for d in dirs:
                third_level_dirs.append(os.path.join(root, d))
        elif current_level > 2:
            break  # No need to go deeper once we reach beyond the third level

    return third_level_dirs

def list_level_two_subdirectories(start_dir):
    level = 0
    two_level_dirs = []

    # Use os.walk to traverse the directory
    for root, dirs, files in os.walk(start_dir):
        # Calculate the current level based on the depth of root
        current_level = root.replace(start_dir, '').count(os.sep)
        
        # Check if it's the third level
        if current_level == 1:
            for d in dirs:
                two_level_dirs.append(os.path.join(root, d))
        elif current_level > 1:
            break  # No need to go deeper once we reach beyond the third level

    return two_level_dirs

# Get the current date
current_date = time.strftime("%Y-%m-%d")
#print("Current date:", current_date)

path='E:\\ppt'
cases_list=os.listdir(path)
slab_1=[]
slab_2=[]
for case in cases_list:
    datetime = os.listdir(path+'\\'+case)[0]
    temp = path+'\\'+case+'\\'+datetime
    if len(os.listdir(temp))==1:
        slab_1.append(case)
    else:
        slab_2.append(case)


#case='ETERNAL_AU090001'
for case in cases_list:
    # Create a presentation object
    prs = Presentation()

    # Add a title slide
    slide_layout = prs.slide_layouts[0]  # 0 is the layout index for a title slide
    slide = prs.slides.add_slide(slide_layout)
# Add title and subtitle to the title slide
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Eternal Trial"
    subtitle.text = "Rapid output"
    

    
    # Get the slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Define the dimensions of the text box
    textbox_width = Inches(4)  # You can adjust this width
    textbox_height = Inches(1)  # You can adjust this height
    
    # Calculate the position to center the text box
    left = (slide_width - textbox_width) / 2
    top = ((slide_height - textbox_height) / 2)+Inches(1/2)
    
    # Add the date text box to the slide
    textbox = slide.shapes.add_textbox(left, top, textbox_width, textbox_height)
    text_frame = textbox.text_frame
    text_frame.text = f"Date: {current_date}"
    # Center the text inside the text box
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    # Add a slide with content
    slide_layout = prs.slide_layouts[5]  # 1 is the layout index for a content slide
    slide = prs.slides.add_slide(slide_layout)
    # Add title and content to the content slide  2
    title = slide.shapes.title
    
    total_cases=220
    title.text = "Summary"
    textbox = slide.shapes.add_textbox(left, top, textbox_width, textbox_height)
    text_frame = textbox.text_frame
    text_frame.text = f"Total cases in Batch : {total_cases}"
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.level = 0  # Set bullet level to 0 to ensure no bullets
    #prs.save('E:\\temp\\presentation.pptx')
    #sys.exit()



    #case='ETERNAL_AU090001'
    if case in slab_1:
        json_file_path = list_level_third_subdirectories(f'E://ppt//{case}')
        if len(json_file_path)==0:
            continue
        json_file = json_file_path[0]+'\\'+'output.json'
        #sys.exit()
        with open(json_file, 'r') as f:
            json_data = json.load(f)
            for k, v in json_data.items():
                if k == 'MeasurementSummary.Thresholding':
                    print(v[0])
            
            
            print(json_data['DICOMHeaderInfo']['Patient']['PatientName'])
            print(json_data['DICOMHeaderInfo']["PerfusionSeries"][0]['StudyDate'])
            print(json_data['DICOMHeaderInfo']["PerfusionSeries"][0]["StudyTime"])
            date=json_data['DICOMHeaderInfo']["PerfusionSeries"][0]['StudyDate']
            time=json_data['DICOMHeaderInfo']["PerfusionSeries"][0]["StudyTime"]
        
        cta_path='E://ppt_data//CTA//BL-SBC_LVO-V1'
        cta_cases_list=os.listdir(cta_path)
        for cta in cta_cases_list:
            if cta==case:
                print(cta)
                #sys.exit()
                
                DateTime= os.listdir(cta_path+'//'+cta)
                #sys.exit()
                series = os.listdir(cta_path+'//'+cta+'//'+DateTime[0])
                png_list=os.listdir(cta_path+'//'+cta+'//'+DateTime[0]+'//'+series[0])
                cta_pngs=[]
                cta_image_path_list=[]
                for file in png_list:
                    s=file.split('.')
                    if s[-1]=='png':
                        cta_pngs.append(file)
                        cta_image_path_list.append(cta_path+'//'+cta+'//'+DateTime[0]+'//'+series[0]+'//'+file)
                # Define image dimensions
                if len(cta_pngs)!=0:
                    
                
                    cta_list_4_images=[]
                    for i in range(1,len(cta_image_path_list)):
                        cta_list_4_images.append(cta_image_path_list[i])
                    #sys.exit()
                    left = Inches(1.5)
                    top = Inches(1.5)
                    height = Inches(2) 
                    width=Inches(3) # Adjust height as needed
                    pic1 = slide.shapes.add_picture(cta_list_4_images[0], left, top, width=width,height=height)
                    
                    # Add the second image
                    
                    left = Inches(5.5)  # Adjust position as needed
                    top = Inches(1.5)
                    height = Inches(2)  # Adjust height as needed
                    pic2 = slide.shapes.add_picture(cta_list_4_images[1], left, top,width=width, height=height)
                    
                    # Add the third image
                    left = Inches(1.5)
                    top = Inches(4.5)
                    height = Inches(2)  # Adjust height as needed
                    pic3 = slide.shapes.add_picture(cta_list_4_images[2], left, top,width=width, height=height)
                    
                    # Add the fourth image
                    
                    left = Inches(5.5)  # Adjust position as needed
                    top = Inches(4.5)
                    height = Inches(2)  # Adjust height as needed
                    pic4 = slide.shapes.add_picture(cta_list_4_images[3], left, top, width=width,height=height)
                    
                    # Add the date text box to the slide
                    # Add a text box at the bottom of the slide
                    left = Inches(2.0)
                    top = Inches(6.5)
                    width = Inches(2)
                    height = Inches(1)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    p = text_frame.add_paragraph()
                    p.text = "Left Hemisphere"
                    p.font.size = Pt(24)  # Set font size
                    
                    left = Inches(6.0)
                    top = Inches(6.5)
                    width = Inches(2)
                    height = Inches(1)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    p = text_frame.add_paragraph()
                    p.text = "Right Hemisphere"
                    p.font.size = Pt(24)  # Set font size
                    #prs.save(f'E:\\temp\\presentation_{case}.pptx')
                    #sys.exit()
                else:
                    continue
            else:
                continue
        
        slide_layout = prs.slide_layouts[5]  # 6
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title 
        for shape in slide.shapes:  ## remove title
            if shape.has_text_frame and shape.text_frame.text == slide.shapes.title.text_frame.text:
                slide.shapes._spTree.remove(shape._element)
        
        # Add the date text box to the slide
        left=1
        top=1
        textbox = slide.shapes.add_textbox(left+Inches(1), top, textbox_width, textbox_height)
        text_frame = textbox.text_frame
        text_frame.text =f"Case {case}_{date}_{time}"
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
        # Specify the path to the image
        png_dir_path=list_level_third_subdirectories('E://ppt_data//montages')
        ncct_png_path=png_dir_path[0]
        ncct_all=os.listdir(ncct_png_path)
        ### volume.txt
        with open('E:\\ppt_data\\volume.txt', 'r') as file:
            for line in file:
                print(line)
                line= line.replace('\n','')
                c=line.split('_')
                CASE=c[0]+'_'+c[1]
                if CASE==case:
                    sss = line.split('\t')
                    vol = sss[-1]
        # Add a text box at the top left of the slide
        left = Inches(6.5)
        top = Inches(0.5)
        width = Inches(4)
        height = Inches(1)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = f"Predicted Volume: {vol}ml"
        p.font.size = Pt(14)  # Set font size
        for ncct in ncct_all:
            s=case.split('_')
            ss=s[-1]
            r=ncct.split('_')
            rr=r[0]
            if rr==ss:
                ncct_path=ncct_png_path+'//'+ncct
        # Add the image to the slide
        left = Inches(1)
        top = Inches(2)
        height = Inches(4)  # Adjust height as needed
        
        slide.shapes.add_picture(ncct_path, left, top, width = Inches(4),height=height)
        
        pred_png_path=png_dir_path[1]
        pred_all=os.listdir(pred_png_path)
        
        for pred in pred_all:
            s=case.split('_')
            ss=s[-1]
            r=pred.split('.')
            rr=r[0]
            if rr==ss:
                pred_path=pred_png_path+'//'+pred
        slide.shapes.add_picture(pred_path, left+Inches(4.5), top, width = Inches(4),height=height)
        
        
        slide_layout = prs.slide_layouts[5]  # new 5 from old 3
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title 
        for shape in slide.shapes:  ## remove title
            if shape.has_text_frame and shape.text_frame.text == slide.shapes.title.text_frame.text:
                slide.shapes._spTree.remove(shape._element)
        textbox_width = Inches(4)  # You can adjust this width
        textbox_height = Inches(1)  # You can adjust this height

        # Calculate the position to center the text box
        left = 1
        top = 1
        json_file_path = list_level_third_subdirectories(f'E://ppt//{case}')
        if len(json_file_path)==0:
            continue
        json_file = json_file_path[0]+'\\'+'output.json'
        #sys.exit()
        with open(json_file, 'r') as f:
            json_data = json.load(f)
            for k, v in json_data.items():
                if k == 'MeasurementSummary.Thresholding':
                    print(v[0])
            
            
            print(json_data['DICOMHeaderInfo']['Patient']['PatientName'])
            print(json_data['DICOMHeaderInfo']["PerfusionSeries"][0]['StudyDate'])
            print(json_data['DICOMHeaderInfo']["PerfusionSeries"][0]["StudyTime"])
            date=json_data['DICOMHeaderInfo']["PerfusionSeries"][0]['StudyDate']
            time=json_data['DICOMHeaderInfo']["PerfusionSeries"][0]["StudyTime"]
            #sys.exit(0)
            print(json_data["MeasurementSummary"])
            thresholds = json_data["MeasurementSummary"]
            print(thresholds["Thresholding"])
            thresholds_list = thresholds["Thresholding"]
            print('\n\nthresholds')
            print(thresholds_list[1]['Parameter'])
            print(thresholds_list[1]['Threshold'])
            print(thresholds_list[1]['Volume'])
            
            print(thresholds_list[5]['Parameter'])
            print(thresholds_list[5]['Volume'])
            #sys.exit()
        Mismatch_volume = abs(thresholds_list[1]['Volume']-thresholds_list[5]['Volume'])
        if thresholds_list[1]['Volume']==0:
            Mismatch_ratio ='Nan'
        else:
            Mismatch_ratio = thresholds_list[5]['Volume']/thresholds_list[1]['Volume']
        # Add the date text box to the slide
        textbox = slide.shapes.add_textbox(left+Inches(1), top, textbox_width, textbox_height)
        text_frame = textbox.text_frame
        text_frame.text =f"Case {case}_{date}_{time}"
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True

        # Specify the path to the image
        image_dir_path=list_level_subdirectories(f'E://ppt//{case}')
        #sys.exit()
        ii=image_dir_path[0].split('\\')
        iii=ii[2]
        match = re.match(r'^\d+', iii)
        iiii=match.group()
        image_path1 = f'{image_dir_path[0]}//series{iiii}01_view01_CBF_Tmax_Mismatch.jpg'  # Replace with your image path
        
        # Add the image to the slide
        image_height = Inches(4.5)
        image_width  = Inches(9)
        left = (slide_width - image_width) / 2
        top = (slide_height - image_height) / 2
        
        try:
            slide.shapes.add_picture(image_path1, left, top+Inches(1/2), width=image_width,height=image_height)
        except:
            continue
        
        slide_layout = prs.slide_layouts[5]  # new 6 from old 4
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title 
        for shape in slide.shapes:  ## remove title
            if shape.has_text_frame and shape.text_frame.text == slide.shapes.title.text_frame.text:
                slide.shapes._spTree.remove(shape._element)
        
        # Add the date text box to the slide
        left=1
        top=1
        textbox = slide.shapes.add_textbox(left+Inches(1), top, textbox_width, textbox_height)
        text_frame = textbox.text_frame
        text_frame.text =f"Case {case}_{date}_{time}"
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
        # Specify the path to the image
        image_dir_path=list_level_subdirectories(f'E://ppt//{case}')
        image_path1 = f'{image_dir_path[0]}//series{iiii}01_view09_AIF_VOF_Curves.jpg'  # Replace with your image path
        image_path2 = f'{image_dir_path[0]}//series{iiii}01_view06_AIF_VOF_Locations.jpg'  # Replace with your image path
        #image_path3 = f'{image_dir_path[0]}//series801_view07_tsMIP_AIF_VOF_Locations.jpg'
        #image_path4 = f'{image_dir_path[0]}//series801_view05_Columned_View.jpg'
        # Save the presentation
        
        # Define image dimensions
        image_width = Inches(4)
        image_height = Inches(4)
        
        # Add the first image
        
        left = Inches(1)
        top = Inches(2)
        height = Inches(4)  # Adjust height as needed
        pic = slide.shapes.add_picture(image_path1, left, top, width = Inches(4),height=height)
        
        # Add the second image
        left = Inches(5.5)  # Adjust position as needed
        pic = slide.shapes.add_picture(image_path2, left, top,width = Inches(4), height=height)
        
        
        
        slide_layout = prs.slide_layouts[5]  # new 7 from old 5
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title 
        for shape in slide.shapes:  ## remove title
            if shape.has_text_frame and shape.text_frame.text == slide.shapes.title.text_frame.text:
                slide.shapes._spTree.remove(shape._element)
        
        # Add the date text box to the slide
        left=1
        top=1
        textbox = slide.shapes.add_textbox(left+Inches(1), top, textbox_width, textbox_height)
        text_frame = textbox.text_frame
        text_frame.text =f"Case {case}_{date}_{time}"
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
        # Specify the path to the image
        #image_dir_path=list_level_subdirectories(f'E://ppt//{case}')
        image_path4 = f'{image_dir_path[0]}//series{iiii}01_view05_Columned_View.jpg'
        left = Inches(3)
        top = Inches(1/2)
        height = Inches(6.5)  # Adjust height as needed
        pic = slide.shapes.add_picture(image_path4, left, top, width = Inches(4),height=height)
        #prs.save('E:\\temp\\presentation.pptx')
        #sys.exit()


        slide_layout = prs.slide_layouts[5]  # 8 for bolus series801_view07_tsMIP_AIF_VOF_Locations
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title 
        for shape in slide.shapes:  ## remove title
            if shape.has_text_frame and shape.text_frame.text == slide.shapes.title.text_frame.text:
                slide.shapes._spTree.remove(shape._element)
        
        # Add the date text box to the slide
        # Add the date text box to the slide
        left=1
        top=1
        textbox = slide.shapes.add_textbox(left+Inches(1), top, textbox_width, textbox_height)
        text_frame = textbox.text_frame
        text_frame.text =f"Case {case}_{date}_{time}"
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
         # Specify the path to the image
        image_dir_path=list_level_subdirectories(f'E://ppt//{case}')
        image_path = f'{image_dir_path[0]}//series{iiii}01_view07_tsMIP_AIF_VOF_Locations.jpg'  # Replace with your image path

         # Define image dimensions
        image_width = Inches(4)
        image_height = Inches(4)
         
         # Add the first image
         
        left = Inches(1)
        top = Inches(2)
        height = Inches(4)  # Adjust height as needed
        image = slide.shapes.add_picture(image_path, left, top, width = Inches(4),height=height)

        # Get image dimensions
        image_width = image.width
        image_height = image.height
        
        # Calculate the position to center the image
        left = (slide_width - image_width) // 2
        top = (slide_height - image_height) // 2
        
        # Move the image to the calculated position
        image.left = left
        image.top = top
        prs.save(f'E:\\temp\\presentation_{case}.pptx')
        #print("YES")
        #sys.exit(0)
    